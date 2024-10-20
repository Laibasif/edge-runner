import os
import io
import PyPDF2
import streamlit as st
from pptx import Presentation
from PIL import Image
from openai import OpenAI

# Initialize session state variables for managing chat history and document index
if 'history' not in st.session_state:
    st.session_state['history'] = []

# Set up the page configuration
st.set_page_config(layout="wide")

# Initialize the Llama model through OpenAI API
@st.cache_resource
def initialize_llama():
    client = OpenAI(
        api_key="YOUR_API_KEY",  # Replace with your OpenAI API key
        base_url="https://api.aimlapi.com",
    )
    return client

# Function to extract text from PDF files
def extract_text_from_pdf(file):
    text = ""
    with io.BytesIO(file.read()) as pdf_file:
        reader = PyPDF2.PdfReader(pdf_file)
        for page in reader.pages:
            text += page.extract_text() + "\n"
    return text

# Function to extract text from PPT files
def extract_text_from_ppt(file):
    text = ""
    presentation = Presentation(file)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to summarize the text using Llama model
def summarize_text(llama_client, text):
    response = llama_client.chat.completions.create(
        model="meta-llama/Llama-3.2-3B-Instruct-Turbo",
        messages=[
            {
                "role": "system",
                "content": "You are an AI assistant who summarizes documents.",
            },
            {
                "role": "user",
                "content": f"Please summarize the following text:\n{text}"
            },
        ],
    )
    return response.choices[0].message.content

def main():
    llama_client = initialize_llama()

    st.title("Llama Model Chatbot")

    # File upload functionality
    uploaded_file = st.file_uploader("Upload a PPT, PDF, or Image file", type=["pptx", "pdf", "png", "jpg", "jpeg"])
    
    if uploaded_file is not None:
        if uploaded_file.type == "application/pdf":
            text = extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            text = extract_text_from_ppt(uploaded_file)
        else:
            # If the uploaded file is an image, we can extract text using OCR (not implemented here)
            st.warning("Image processing is not implemented yet. Please upload a PDF or PPT.")
            text = ""
        
        if text:
            summary = summarize_text(llama_client, text)
            st.markdown("### Summary:")
            st.write(summary)

    # Chat interface
    chat_container = st.container()
    with chat_container:
        for message in st.session_state['history']:
            with st.chat_message(message["role"]):
                st.markdown(message["content"])

    user_input = st.chat_input("Enter your query:")

    if user_input:
        with st.chat_message("user"):
            st.markdown(user_input)
        st.session_state['history'].append({"role": "user", "content": user_input})

        with st.chat_message("assistant"):
            with st.spinner("Generating response..."):
                # Generate response using OpenAI Llama model
                response = llama_client.chat.completions.create(
                    model="meta-llama/Llama-3.2-3B-Instruct-Turbo",
                    messages=[
                        {
                            "role": "system",
                            "content": "You are an AI assistant who knows everything.",
                        },
                        {
                            "role": "user",
                            "content": user_input
                        },
                    ],
                )
                full_response = response.choices[0].message.content
                st.markdown(full_response)
        st.session_state['history'].append({"role": "assistant", "content": full_response})
        st.rerun()  # Force a rerun to update the chat display

    # Add a clear button
    if st.button("Clear Chat"):
        st.session_state['history'] = []
        st.rerun()

if __name__ == "__main__":
    main()
